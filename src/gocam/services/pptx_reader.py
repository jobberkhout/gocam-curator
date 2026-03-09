"""Extract text and embedded images from PowerPoint files."""

from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path


@dataclass
class SlideContent:
    slide_number: int
    text: str           # all text from shapes, title, body
    notes: str          # speaker notes
    images: list[bytes] = field(default_factory=list)  # raw bytes of embedded images


def read_pptx(path: Path) -> list[SlideContent]:
    """Extract all slides from a .pptx file.

    Returns one SlideContent per slide with text, notes, and embedded images.
    Note: this extracts content from slides, not a rendered pixel image of the slide.
    Full slide rendering requires LibreOffice (not bundled).
    """
    from pptx import Presentation
    from pptx.shapes.picture import Picture

    prs = Presentation(path)
    slides: list[SlideContent] = []

    for i, slide in enumerate(prs.slides):
        # Collect all text from every shape (title, body, text boxes, tables)
        text_lines: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        text_lines.append(line)
            # Extract table cell text
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text:
                            text_lines.append(cell_text)

        # Speaker notes
        notes = ""
        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            notes = notes_tf.text.strip()

        # Embedded images (Picture shapes)
        images: list[bytes] = []
        for shape in slide.shapes:
            if isinstance(shape, Picture):
                try:
                    images.append(shape.image.blob)
                except Exception:
                    pass

        slides.append(SlideContent(
            slide_number=i + 1,
            text="\n".join(text_lines),
            notes=notes,
            images=images,
        ))

    return slides
